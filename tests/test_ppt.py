import io
from pptx import Presentation
from extractors.ppt import extract_ppt


def _make_ppt_with_text(text: str) -> bytes:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = text
    body = slide.placeholders[1]
    body.text = "Coverage: ₹5,00,000"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_ppt_returns_text_and_images():
    ppt_bytes = _make_ppt_with_text("HealthFlex Bronze Plan")
    text, images = extract_ppt(ppt_bytes)
    assert isinstance(text, str)
    assert isinstance(images, list)


def test_extract_ppt_text_contains_title():
    ppt_bytes = _make_ppt_with_text("HealthFlex Bronze Plan")
    text, _ = extract_ppt(ppt_bytes)
    assert "HealthFlex Bronze Plan" in text


def test_extract_ppt_text_contains_body():
    ppt_bytes = _make_ppt_with_text("My Slide")
    text, _ = extract_ppt(ppt_bytes)
    assert "Coverage" in text
