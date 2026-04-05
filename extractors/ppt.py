import io
from pptx import Presentation


def extract_ppt(file_bytes: bytes) -> tuple[str, list[bytes]]:
    """Extract all text and embedded images from a PPT file."""
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    images = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        text_parts.append(line)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_bytes = shape.image.blob
                images.append(image_bytes)

    return "\n".join(text_parts), images
